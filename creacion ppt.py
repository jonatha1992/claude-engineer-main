from pptx import Presentation
from pptx.util import Inches

# Crear una nueva presentación
prs = Presentation()

# Diapositiva 1: Título y Introducción
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Optimización de Transcripción de Audio"
subtitle.text = "Desarrollo de un programa para la conversión y transcripción automatizada"
content = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
text_frame = content.text_frame
text_frame.text = "Este proyecto aborda la necesidad de mejorar la precisión y eficiencia en la transcripción de audio, "
p = text_frame.add_paragraph()
p.text = "proporcionando una solución que automatiza la conversión y procesamiento de archivos de audio."
p.space_after = Inches(0.1)

# Diapositiva 2: Conversión de Audio
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Conversión de Audio"
content.text = (
    "• Conversión de archivos de audio a formato WMA, que es más compatible con el sistema de transcripción.\n"
    "• Conversión a mono y reducción de interferencia para asegurar que la voz sea clara y dominante en la grabación.\n"
    "• Segmentación en partes de 10 segundos para mejorar precisión y facilitar la transcripción de audios largos."
)
extra_content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
extra_frame = extra_content.text_frame
extra_frame.text = "Estas técnicas permiten que el programa maneje eficientemente el procesamiento de audio, incluso en condiciones subóptimas."

# Diapositiva 3: Transcripción y Traducción
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Transcripción y Traducción"
content.text = (
    "• El programa convierte el audio en texto utilizando las potentes librerías de Google Translator.\n"
    "• Configuración flexible que permite establecer tanto el idioma de entrada como el de salida.\n"
    "• Ejemplo práctico: Transcripción de audio en inglés y su traducción automática al español con alta precisión."
)
extra_content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
extra_frame = extra_content.text_frame
extra_frame.text = "La integración de Google Translator añade una capa adicional de funcionalidad, permitiendo una transcripción multilingüe."

# Diapositiva 4: Revisión y Exportación de Texto
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Revisión y Exportación de Texto"
content.text = (
    "• El texto transcrito incluye marcas de tiempo, lo que facilita la revisión y corrección.\n"
    "• El usuario puede reproducir el audio correspondiente desde el punto exacto para verificar la precisión de la transcripción.\n"
    "• El texto puede exportarse fácilmente a un archivo .txt, permitiendo su manipulación externa en otras aplicaciones."
)
extra_content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
extra_frame = extra_content.text_frame
extra_frame.text = "Estas características aseguran que el usuario tenga un control total sobre el proceso de transcripción y su resultado final."

# Diapositiva 5: Eficiencia del Programa
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Eficiencia del Programa"
content.text = (
    "• Ahorro significativo de tiempo en comparación con la transcripción manual, que puede ser tediosa y propensa a errores.\n"
    "• Ejemplo: Una transcripción de 7 minutos realizada por el programa toma solo una fracción del tiempo que requeriría hacerlo a mano.\n"
    "• El rendimiento del programa puede variar dependiendo del hardware de la PC, ya que utiliza intensivamente la memoria y el procesamiento."
)
extra_content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
extra_frame = extra_content.text_frame
extra_frame.text = "El uso del programa no solo ahorra tiempo, sino que también mejora la precisión, eliminando el esfuerzo manual repetitivo."

# Diapositiva 6: Conclusiones
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Conclusiones"
content.text = (
    "• El programa ofrece beneficios claros en términos de precisión y eficiencia en la transcripción de audio.\n"
    "• La capacidad de traducir automáticamente el texto transcrito agrega un valor significativo para usuarios multilingües.\n"
    "• Esta herramienta representa un avance en la automatización de procesos, mejorando la productividad y reduciendo el tiempo requerido para tareas repetitivas."
)
extra_content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
extra_frame = extra_content.text_frame
extra_frame.text = "Con esta solución, los usuarios pueden enfocarse en el contenido, dejando el trabajo pesado de la transcripción al programa."

# Guardar la presentación
prs.save('TuPresentacionRellena.pptx')

print("Presentación creada y guardada exitosamente como 'TuPresentacionRellena.pptx'")
